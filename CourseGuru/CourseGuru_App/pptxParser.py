import nltk
import re

from pptx import Presentation

from nltk.corpus import stopwords
from nltk.tokenize.moses import MosesDetokenizer

from CourseGuru_App.models import botanswers

#testing PRP
#def restructForDB (text): 
def restructForDB (text, cid, catID, fileid): 
    detokenizer = MosesDetokenizer() 
    
    botSearch = text.replace('<br>', '')
    rgx = re.compile('[^a-zA-Z0-9 \n\.]')
    
    botSearch = rgx.sub('', botSearch)
    searchList = nltk.word_tokenize(botSearch, 'english')
    
    botSearch = [word for word in searchList if word not in stopwords.words('english')]
    
    detokenizer.detokenize(botSearch, return_str=True)
    botSearch = ' '.join(botSearch)
    botSearch = botSearch.lower()
    
    botanswers.objects.create(answer = text, rating = 0, category_id = catID.id, entities = botSearch, course_id = cid, file_id = fileid)
#testing PRP    
#    print(botSearch)


def parsePPTX(file, cid, catID, fileid, docType):
    parseFile = Presentation(file)
    dataArray = [] 

    for slide in parseFile.slides:
        dataArray.append('')
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.is_placeholder == True and 'FOOTER' not in str(shape.placeholder_format.type) or shape.is_placeholder == False:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text != '':
                            text = paragraph.text
                            if re.match('^[0-9]*$', text) == None:
                                if text == '.':
                                    dataArray[-1] += '' + (text)
                                else: 
                                    dataArray[-1] += '' + (text) + '<br>'                            
            elif shape.has_table:
                table = shape.table
                
                numCols = len(table.columns)
                numRows = len(table.rows)           
                
                i = 0 
                text = ''
                while i <numRows:
                    j=0
                    while j<numCols:
                        check = (table.cell(i, j)).text_frame.text
                        if check != ' ' and check != '': 
                            if text == '':
                                text = '<br>' + check
                            elif j > 0:
                                text = text + ' -- ' + check
                            else:
                                text = text + check
                        if j == numCols-1:
                            text = text + '<br>'
                        j+=1
                    i+=1
                dataArray[-1] += ' ' + (text) 
                
    file.close()
    for n in dataArray:
        restructForDB(n, cid, catID, fileid)
#testing PRP   
#        print(n)